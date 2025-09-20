"""Text extraction utilities for various document formats."""

import os
import re
from typing import List, Tuple, Optional
from pptx import Presentation
from dataclasses import dataclass
import PyPDF2
import pdfplumber
import docx
import markdown
from config import SUPPORTED_EXTENSIONS, OCR_ENABLED


@dataclass
class ExtractedChunk:
    """Represents a chunk of extracted text with metadata."""
    text: str
    filename: str
    page: Optional[int] = None
    section: Optional[str] = None
    char_offset: int = 0
    course: Optional[str] = None


class TextExtractor:
    """Handles text extraction from various document formats."""
    
    def __init__(self):
        self.extractors = {
            '.pdf': self._extract_pdf,
            '.md': self._extract_markdown,
            '.txt': self._extract_text,
            '.docx': self._extract_docx,
            '.pptx': self._extract_pptx
        }
    
    def extract_from_file(self, filepath: str, course: str = None) -> List[ExtractedChunk]:
        """Extract text from a file and return chunks with metadata."""
        if not os.path.exists(filepath):
            raise FileNotFoundError(f"File not found: {filepath}")
        
        ext = os.path.splitext(filepath)[1].lower()
        if ext not in SUPPORTED_EXTENSIONS:
            raise ValueError(f"Unsupported file extension: {ext}")
        
        filename = os.path.basename(filepath)
        extractor = self.extractors.get(ext)
        
        if not extractor:
            raise ValueError(f"No extractor available for {ext}")
        
        try:
            chunks = extractor(filepath, filename, course)
            return self._clean_chunks(chunks)
        except Exception as e:
            print(f"Error extracting from {filepath}: {str(e)}")
            return []
    
    def _extract_pdf(self, filepath: str, filename: str, course: str) -> List[ExtractedChunk]:
        """Extract text from PDF file with page numbers using pdfplumber for better quality."""
        chunks = []
        
        try:
            # Try pdfplumber first for better text extraction
            with pdfplumber.open(filepath) as pdf:
                for page_num, page in enumerate(pdf.pages, 1):
                    try:
                        text = page.extract_text()
                        if text and text.strip():
                            # Clean and process the text
                            cleaned_text = self._remove_headers_footers(text)
                            
                            # Create chunk with page metadata
                            chunk = ExtractedChunk(
                                text=cleaned_text,
                                filename=filename,
                                page=page_num,
                                course=course,
                                char_offset=0  # Will be updated during chunking
                            )
                            chunks.append(chunk)
                            
                    except Exception as e:
                        print(f"Error extracting page {page_num} from {filepath}: {e}")
                        continue
                    
        except Exception as e:
            print(f"Error reading PDF with pdfplumber {filepath}: {e}")
            # Fallback to PyPDF2 if pdfplumber fails
            try:
                with open(filepath, 'rb') as file:
                    pdf_reader = PyPDF2.PdfReader(file)
                    
                    for page_num, page in enumerate(pdf_reader.pages, 1):
                        try:
                            text = page.extract_text()
                            if text.strip():
                                # Clean and process the text
                                cleaned_text = self._remove_headers_footers(text)
                                
                                # Create chunk with page metadata
                                chunk = ExtractedChunk(
                                    text=cleaned_text,
                                    filename=filename,
                                    page=page_num,
                                    course=course,
                                    char_offset=0  # Will be updated during chunking
                                )
                                chunks.append(chunk)
                                
                        except Exception as e:
                            print(f"Error extracting page {page_num} from {filepath}: {e}")
                            continue
                        
            except Exception as e2:
                print(f"Error reading PDF with PyPDF2 {filepath}: {e2}")
                return []
    
        return chunks
    
    def _extract_markdown(self, filepath: str, filename: str, course: str) -> List[ExtractedChunk]:
        """Extract text from Markdown with section preservation."""
        with open(filepath, 'r', encoding='utf-8') as file:
            content = file.read()
        
        chunks = []
        current_section = None
        current_text = []
        
        lines = content.split('\n')
        for line in lines:
            # Check for headers
            header_match = re.match(r'^(#{1,6})\s+(.+)$', line)
            if header_match:
                # Save previous section if exists
                if current_text:
                    chunks.append(ExtractedChunk(
                        text='\n'.join(current_text),
                        filename=filename,
                        section=current_section,
                        course=course
                    ))
                    current_text = []
                
                current_section = header_match.group(2)
                current_text.append(line)
            else:
                current_text.append(line)
        
        # Add final section
        if current_text:
            chunks.append(ExtractedChunk(
                text='\n'.join(current_text),
                filename=filename,
                section=current_section,
                course=course
            ))
        
        return chunks
    
    def _extract_text(self, filepath: str, filename: str, course: str) -> List[ExtractedChunk]:
        """Extract text from plain text files."""
        with open(filepath, 'r', encoding='utf-8') as file:
            content = file.read()
        
        return [ExtractedChunk(
            text=content,
            filename=filename,
            course=course
        )]
    
    def _extract_docx(self, filepath: str, filename: str, course: str) -> List[ExtractedChunk]:
        """Extract text from Word documents."""
        doc = docx.Document(filepath)
        chunks = []
        current_section = None
        
        for para in doc.paragraphs:
            text = para.text.strip()
            if not text:
                continue
            
            # Check if this looks like a heading
            if para.style.name.startswith('Heading') or len(text) < 100 and text.isupper():
                current_section = text
            
            chunks.append(ExtractedChunk(
                text=text,
                filename=filename,
                section=current_section,
                course=course
            ))
        
        return chunks
    
    def _clean_chunks(self, chunks: List[ExtractedChunk]) -> List[ExtractedChunk]:
        """Clean and normalize extracted text chunks."""
        cleaned_chunks = []
        
        for chunk in chunks:
            # Remove excessive whitespace
            text = re.sub(r'\s+', ' ', chunk.text.strip())
            
            # Skip very short chunks
            if len(text) < 50:
                continue
            
            # Remove common headers/footers patterns
            text = self._remove_headers_footers(text)
            
            if text.strip():
                chunk.text = text
                cleaned_chunks.append(chunk)
        
        return cleaned_chunks
    
    def _extract_pptx(self, filepath: str, filename: str, course: str) -> List[ExtractedChunk]:
        """Extract text from PowerPoint presentation."""
        chunks = []
        
        try:
            prs = Presentation(filepath)
            
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_text = []
                
                # Extract text from all shapes in the slide
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        slide_text.append(shape.text.strip())
                
                if slide_text:
                    # Combine all text from the slide
                    combined_text = '\n'.join(slide_text)
                    cleaned_text = self._remove_headers_footers(combined_text)
                    
                    # Create chunk with slide metadata
                    chunk = ExtractedChunk(
                        text=cleaned_text,
                        filename=filename,
                        page=slide_num,  # Use slide number as page
                        course=course,
                        char_offset=0  # Will be updated during chunking
                    )
                    chunks.append(chunk)
                    
        except Exception as e:
            print(f"Error extracting PPTX {filepath}: {e}")
            return []
        
        return chunks
    
    def _remove_headers_footers(self, text: str) -> str:
        """Remove common header/footer patterns and fix OCR artifacts."""
        # Remove running headers first
        text = re.sub(r'Unit \d+:\s*\w+\s*\d*', '', text)
        text = re.sub(r'Chapter \d+\s*-\s*[\w\s]+', '', text)
        
        # De-hyphenate words split across lines
        text = re.sub(r'(\w+)-\s*\n\s*(\w+)', r'\1\2', text)
        text = re.sub(r'(\w+)-\s+(\w+)', r'\1\2', text)
        
        # Aggressive OCR artifact fixing
        ocr_fixes = {
            r'\bapprov\s+ed\b': 'approved',
            r'\bforanumber\b': 'for a number',
            r'\bstem\s+Maintenance\b': 'System Maintenance',
            r'\bmaintenanceofthe\b': 'maintenance of the',
            r'\btorun\b': 'to run',
            r'\bsmoothlyasw\s+ell\b': 'smoothly as well',
            r'\basreduce\b': 'as reduce',
            r'\briskofbreak\s+downs\b': 'risk of breakdowns',
            r'\bhavingawell-designed\b': 'having a well-designed',
            r'\bhavingawell\b': 'having a well',
            r'\bisafrequently\b': 'is a frequently',
            r'\bassistsyouin\b': 'assists you in',
            r'\bmaintenanceandupgrad\s+es\b': 'maintenance and upgrades',
            r'\bareequally\b': 'are equally',
            r'\btoserversisoften\b': 'to servers is often',
            r'\brity\s+features\b': 'security features',
            r'\bcurity\s+features\b': 'security features',
            r'\ban\s+d\b': 'and',
            r'\bm\s+eans\b': 'means',
            r'\bf\s+eatures\b': 'features',
            r'\bw\s+ell\b': 'well',
            r'\bth\s+e\b': 'the',
            r'\bof\s+the\b': 'of the',
            r'\bin\s+the\b': 'in the',
            r'\bmanag\s+ement\b': 'management',
            r'\bsched\s+ule\b': 'schedule',
            r'\bsecur\s+ity\b': 'security',
            r'\bsyst\s+em\b': 'system'
        }
        
        for pattern, replacement in ocr_fixes.items():
            text = re.sub(pattern, replacement, text, flags=re.IGNORECASE)
        
        # General pattern: fix single letters separated by spaces
        text = re.sub(r'\b([a-z])\s+([a-z])\s+([a-z]+)\b', r'\1\2\3', text)
        text = re.sub(r'\b([a-z]+)\s+([a-z])\s+([a-z]+)\b', r'\1\2\3', text)
        
        # Remove unicode bullets and special characters that cause encoding issues
        text = re.sub(r'[\u25aa\u2022\u2023\u25cf\u25e6\u2043\u204c\u204d\u2219\u25b8]', '•', text)  # Replace bullets
        text = re.sub(r'[^\x00-\x7F]+', ' ', text)  # Remove other non-ASCII characters
        
        # Remove page numbers at start/end
        text = re.sub(r'^\d+\s*', '', text)
        text = re.sub(r'\s*\d+$', '', text)
        
        # Remove duplicate headers/footers
        text = re.sub(r'(Unit \d+:.*?)\1+', r'\1', text)
        
        # Remove common footer patterns
        footer_patterns = [
            r'Page \d+ of \d+',
            r'Copyright.*',
            r'All rights reserved.*',
            r'Confidential.*',
            r'^\s*\d+\s*$'
        ]
        
        for pattern in footer_patterns:
            text = re.sub(pattern, '', text, flags=re.IGNORECASE | re.MULTILINE)
        
        # Clean up excessive whitespace
        text = re.sub(r'\s+', ' ', text)
        
        return text.strip()


def extract_text_from_directory(directory: str, course: str = None) -> List[ExtractedChunk]:
    """Extract text from all supported files in a directory."""
    extractor = TextExtractor()
    all_chunks = []
    
    for root, dirs, files in os.walk(directory):
        for file in files:
            if any(file.lower().endswith(ext) for ext in SUPPORTED_EXTENSIONS):
                filepath = os.path.join(root, file)
                try:
                    chunks = extractor.extract_from_file(filepath, course)
                    all_chunks.extend(chunks)
                    print(f"Extracted {len(chunks)} chunks from {file}")
                except Exception as e:
                    print(f"Failed to extract from {file}: {str(e)}")
    
    return all_chunks
